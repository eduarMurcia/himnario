import sys
import tempfile
import unittest
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from hymn_studio.models import ImageOverlay, Lyrics, SlideTemplate, Stanza, TextOverlay
from hymn_studio.services.slide_generator import PillowSlideBuilder, PptxSlideBuilder

WINDOWS_ARIAL = Path("C:/Windows/Fonts/arial.ttf")


@unittest.skipUnless(WINDOWS_ARIAL.exists(), "Requires a system TTF font (arial.ttf).")
class PillowSlideBuilderTest(unittest.TestCase):
    def test_builds_one_png_per_stanza_plus_cover(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            background = self._make_background(temp_path / "bg.png")

            template = SlideTemplate(
                background=background,
                font_path=WINDOWS_ARIAL,
                font_size=60,
                text_box=(20, 20, 360, 260),
            )
            lyrics = Lyrics(
                title="Test Hymn",
                stanzas=[
                    Stanza(text="First verse line one\nSecond line"),
                    Stanza(text="Chorus line", is_chorus=True),
                ],
            )

            output_dir = temp_path / "slides"
            paths = PillowSlideBuilder(template, template).build(lyrics, output_dir)

            self.assertEqual(len(paths), 3)
            for path in paths:
                self.assertTrue(path.exists())
                with Image.open(path) as image:
                    self.assertEqual(image.size, (400, 300))

    def test_renders_extra_overlays_and_image_overlays(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            background = self._make_background(temp_path / "bg.png")
            ornament = self._make_ornament(temp_path / "ornament.png")

            cover_template = SlideTemplate(
                background=background,
                font_path=WINDOWS_ARIAL,
                font_size=60,
                text_box=(20, 20, 360, 60),
                align="left",
                wrap=False,
                extra_overlays=(
                    TextOverlay(
                        text="Subtitle",
                        font_path=WINDOWS_ARIAL,
                        font_size=20,
                        text_box=(20, 100, 200, 30),
                        align="left",
                        wrap=False,
                    ),
                ),
                image_overlays=(ImageOverlay(image=ornament, box=(20, 140, 100, 10)),),
            )
            stanza_template = SlideTemplate(
                background=background, font_path=WINDOWS_ARIAL, font_size=40
            )
            lyrics = Lyrics(title="Cover Title", stanzas=[Stanza(text="Verse")])

            output_dir = temp_path / "slides"
            paths = PillowSlideBuilder(cover_template, stanza_template).build(lyrics, output_dir)

            self.assertEqual(len(paths), 2)
            with Image.open(paths[0]) as cover_image:
                self.assertEqual(cover_image.mode, "RGBA")
                self.assertEqual(cover_image.size, (400, 300))

    def _make_background(self, path: Path) -> Path:
        Image.new("RGB", (400, 300), color=(10, 20, 30)).save(path)
        return path

    def _make_ornament(self, path: Path) -> Path:
        Image.new("RGBA", (200, 20), color=(246, 235, 213, 255)).save(path)
        return path


class PptxSlideBuilderTest(unittest.TestCase):
    def test_fills_cover_title_and_duplicates_stanza_slide(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            cover_path = self._make_template(temp_path / "cover.pptx", "COVER PLACEHOLDER")
            stanza_path = self._make_template(temp_path / "stanza.pptx", "STANZA PLACEHOLDER")

            lyrics = Lyrics(
                title="Test Hymn",
                stanzas=[Stanza(text="Verse one"), Stanza(text="Verse two", is_chorus=True)],
            )

            output_dir = temp_path / "out"
            combined_path = PptxSlideBuilder(cover_path, stanza_path).build(lyrics, output_dir)

            self.assertTrue(combined_path.exists())
            presentation = Presentation(str(combined_path))
            self.assertEqual(len(list(presentation.slides)), 3)

            texts = [
                shape.text_frame.text
                for slide in presentation.slides
                for shape in slide.shapes
                if shape.has_text_frame
            ]
            self.assertIn("Test Hymn", texts)
            self.assertIn("Verse one", texts)
            self.assertIn("Verse two", texts)

    def test_preserves_original_font_when_filling_text(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            cover_path = self._make_styled_template(temp_path / "cover.pptx", "COVER")
            stanza_path = self._make_styled_template(
                temp_path / "stanza.pptx", "Line one\nLine two"
            )

            lyrics = Lyrics(title="Styled Title", stanzas=[Stanza(text="A\nB\nC")])

            combined_path = PptxSlideBuilder(cover_path, stanza_path).build(
                lyrics, temp_path / "out"
            )

            presentation = Presentation(str(combined_path))
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            self.assertEqual(run.font.name, "Garamond")
                            self.assertEqual(run.font.size, Pt(90))
                            self.assertTrue(run.font.bold)
                            self.assertEqual(run.font.color.rgb, RGBColor(0x6B, 0x1E, 0x2E))

    def _make_styled_template(self, path: Path, placeholder_text: str) -> Path:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        tf = textbox.text_frame
        tf.text = placeholder_text
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Garamond"
                run.font.size = Pt(90)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x6B, 0x1E, 0x2E)
        presentation.save(str(path))
        return path

    def test_duplicates_stanza_slide_with_background_picture(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            cover_path = self._make_template(temp_path / "cover.pptx", "COVER PLACEHOLDER")
            stanza_path = self._make_template_with_picture(
                temp_path, temp_path / "stanza.pptx", "STANZA PLACEHOLDER"
            )

            lyrics = Lyrics(
                title="Test Hymn",
                stanzas=[Stanza(text="Verse one"), Stanza(text="Verse two")],
            )

            output_dir = temp_path / "out"
            combined_path = PptxSlideBuilder(cover_path, stanza_path).build(lyrics, output_dir)

            presentation = Presentation(str(combined_path))
            self.assertEqual(len(list(presentation.slides)), 3)

            # Each duplicated slide must keep a working, independent image relationship.
            for slide in list(presentation.slides)[1:]:
                pictures = [shape for shape in slide.shapes if shape.shape_type == 13]
                self.assertEqual(len(pictures), 1)
                self.assertTrue(pictures[0].image.blob)

    def _make_template_with_picture(
        self, temp_path: Path, path: Path, placeholder_text: str
    ) -> Path:
        image_path = temp_path / f"{path.stem}_bg.png"
        Image.new("RGB", (40, 30), color=(10, 20, 30)).save(image_path)

        presentation = Presentation()
        blank_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(image_path), Inches(0), Inches(0), Inches(4), Inches(3))
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        textbox.text_frame.text = placeholder_text
        presentation.save(str(path))
        return path

    def _make_template(self, path: Path, placeholder_text: str) -> Path:
        presentation = Presentation()
        blank_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        textbox.text_frame.text = placeholder_text
        presentation.save(str(path))
        return path

    def test_fills_text_nested_inside_groups(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            cover_path = self._make_nested_group_template(temp_path / "cover.pptx")
            stanza_path = self._make_template(temp_path / "stanza.pptx", "STANZA PLACEHOLDER")

            lyrics = Lyrics(title="Grouped Title", stanzas=[Stanza(text="Verse one")])

            combined_path = PptxSlideBuilder(cover_path, stanza_path).build(
                lyrics, temp_path / "out"
            )

            presentation = Presentation(str(combined_path))
            texts = self._collect_texts(presentation.slides[0].shapes)
            self.assertIn("Grouped Title", texts)
            self.assertIn("Static Subtitle", texts)

    def _collect_texts(self, shapes) -> list[str]:
        texts: list[str] = []
        for shape in shapes:
            if shape.shape_type == 6:
                texts.extend(self._collect_texts(shape.shapes))
            elif shape.has_text_frame:
                texts.append(shape.text_frame.text)
        return texts

    def _make_nested_group_template(self, path: Path) -> Path:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        title_box.text_frame.text = "TITLE PLACEHOLDER"
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        subtitle_box.text_frame.text = "Static Subtitle"

        inner_group = slide.shapes.add_group_shape([title_box, subtitle_box])
        slide.shapes.add_group_shape([inner_group])

        presentation.save(str(path))
        return path


if __name__ == "__main__":
    unittest.main()
