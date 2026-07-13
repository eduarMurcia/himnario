from __future__ import annotations

import copy
import io
import shutil
import tempfile
from pathlib import Path
from typing import Protocol

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.slide import Slide

from hymn_studio.models import ImageOverlay, Lyrics, SlideTemplate, TextOverlay


class SlideGenerationError(RuntimeError):
    pass


class SlideGenerator(Protocol):
    """Builds slide PNGs for a hymn's lyrics."""

    def build(self, lyrics: Lyrics, output_dir: Path) -> list[Path]: ...


class PillowSlideBuilder:
    """Renders hymn slides by drawing text onto template backgrounds with Pillow.

    Prototype engine kept independent of PowerPoint/COM automation, so it can run
    and be tested anywhere Pillow is installed.
    """

    def __init__(self, cover_template: SlideTemplate, stanza_template: SlideTemplate) -> None:
        self._cover_template = cover_template
        self._stanza_template = stanza_template

    def build(self, lyrics: Lyrics, output_dir: Path) -> list[Path]:
        if not lyrics.stanzas:
            raise SlideGenerationError("Lyrics must have at least one stanza.")

        output_dir.mkdir(parents=True, exist_ok=True)
        paths: list[Path] = []

        cover_path = output_dir / "0000_cover.png"
        self._render_slide(self._cover_template, lyrics.title, cover_path)
        paths.append(cover_path)

        for index, stanza in enumerate(lyrics.stanzas, start=1):
            slide_path = output_dir / f"{index:04d}_slide.png"
            self._render_slide(self._stanza_template, stanza.text, slide_path)
            paths.append(slide_path)

        return paths

    def _render_slide(self, template: SlideTemplate, text: str, output_path: Path) -> None:
        with Image.open(template.background) as background:
            image = background.convert("RGBA").copy()

        for overlay in template.image_overlays:
            self._paste_image_overlay(image, overlay)

        draw = ImageDraw.Draw(image)
        self._draw_text_block(
            draw,
            text=text,
            font_path=template.font_path,
            font_size=template.font_size,
            text_box=template.text_box,
            text_color=template.text_color,
            align=template.align,
            line_spacing=template.line_spacing,
            stroke_width=template.stroke_width,
            stroke_color=template.stroke_color,
            fit_to_box=True,
            wrap=template.wrap,
        )

        for overlay in template.extra_overlays:
            self._draw_overlay(draw, overlay)

        image.save(output_path, "PNG")

    def _paste_image_overlay(self, image: Image.Image, overlay: ImageOverlay) -> None:
        box_x, box_y, box_width, box_height = overlay.box
        with Image.open(overlay.image) as ornament:
            resized = ornament.convert("RGBA").resize((int(box_width), int(box_height)))
        image.paste(resized, (int(box_x), int(box_y)), resized)

    def _draw_overlay(self, draw: ImageDraw.ImageDraw, overlay: TextOverlay) -> None:
        self._draw_text_block(
            draw,
            text=overlay.text,
            font_path=overlay.font_path,
            font_size=overlay.font_size,
            text_box=overlay.text_box,
            text_color=overlay.text_color,
            align=overlay.align,
            line_spacing=overlay.line_spacing,
            stroke_width=0,
            stroke_color=(0, 0, 0),
            fit_to_box=False,
            wrap=overlay.wrap,
        )

    def _draw_text_block(
        self,
        draw: ImageDraw.ImageDraw,
        *,
        text: str,
        font_path: Path,
        font_size: int,
        text_box: tuple[int, int, int, int],
        text_color: tuple[int, int, int],
        align: str,
        line_spacing: float,
        stroke_width: int,
        stroke_color: tuple[int, int, int],
        fit_to_box: bool,
        wrap: bool = True,
    ) -> None:
        box_x, box_y, box_width, box_height = text_box

        if not wrap:
            # Mirrors PowerPoint's wrap="none" + spAutoFit boxes: single line,
            # left-anchored at the box, vertically centered in the box height.
            font = ImageFont.truetype(str(font_path), font_size)
            bbox = draw.textbbox((0, 0), text, font=font)
            text_height = bbox[3] - bbox[1]
            y = box_y + max(box_height - text_height, 0) / 2
            draw.text((box_x, y), text, font=font, fill=text_color)
            return

        if fit_to_box:
            resolved_size, lines = self._fit_text(
                draw, font_path, font_size, line_spacing, text, box_width, box_height
            )
        else:
            resolved_size = font_size
            font = ImageFont.truetype(str(font_path), resolved_size)
            lines = self._wrap_text(draw, font, text, box_width)

        font = ImageFont.truetype(str(font_path), resolved_size)
        line_height = resolved_size * line_spacing
        total_height = line_height * len(lines)
        start_y = box_y + max(box_height - total_height, 0) / 2

        for offset, line in enumerate(lines):
            bbox = draw.textbbox((0, 0), line, font=font, stroke_width=stroke_width)
            line_width = bbox[2] - bbox[0]
            if align == "center":
                x = box_x + max(box_width - line_width, 0) / 2
            else:
                x = box_x
            y = start_y + offset * line_height
            draw.text(
                (x, y),
                line,
                font=font,
                fill=text_color,
                stroke_width=stroke_width,
                stroke_fill=stroke_color,
            )

    def _fit_text(
        self,
        draw: ImageDraw.ImageDraw,
        font_path: Path,
        font_size: int,
        line_spacing: float,
        text: str,
        box_width: float,
        box_height: float,
    ) -> tuple[int, list[str]]:
        while font_size > 8:
            font = ImageFont.truetype(str(font_path), font_size)
            lines = self._wrap_text(draw, font, text, box_width)
            line_height = font_size * line_spacing
            total_height = line_height * len(lines)
            if total_height <= box_height:
                return font_size, lines
            font_size -= 2

        font = ImageFont.truetype(str(font_path), font_size)
        return font_size, self._wrap_text(draw, font, text, box_width)

    def _wrap_text(
        self,
        draw: ImageDraw.ImageDraw,
        font: ImageFont.FreeTypeFont,
        text: str,
        box_width: float,
    ) -> list[str]:
        lines: list[str] = []
        for raw_line in text.splitlines() or [""]:
            words = raw_line.split()
            if not words:
                lines.append("")
                continue

            current = words[0]
            for word in words[1:]:
                candidate = f"{current} {word}"
                bbox = draw.textbbox((0, 0), candidate, font=font)
                if bbox[2] - bbox[0] <= box_width:
                    current = candidate
                else:
                    lines.append(current)
                    current = word
            lines.append(current)

        return lines


class PptxSlideBuilder:
    """Builds a combined .pptx by filling the existing cover and stanza templates.

    Assumes each template is a single-slide .pptx with one text placeholder to fill
    (the title on the cover template, the lyrics on the stanza template). The stanza
    slide is duplicated once per stanza and its placeholder text replaced.
    """

    def __init__(self, cover_template: Path, stanza_template: Path) -> None:
        self._cover_template = cover_template
        self._stanza_template = stanza_template

    def build(self, lyrics: Lyrics, output_dir: Path) -> Path:
        if not lyrics.stanzas:
            raise SlideGenerationError("Lyrics must have at least one stanza.")

        output_dir.mkdir(parents=True, exist_ok=True)
        combined_path = output_dir / "combined.pptx"
        shutil.copyfile(self._cover_template, combined_path)

        presentation = Presentation(str(combined_path))
        if not presentation.slides:
            raise SlideGenerationError(f"Cover template has no slides: {self._cover_template}")

        self._fill_placeholder_text(presentation.slides[0], lyrics.title)

        stanza_source = Presentation(str(self._stanza_template))
        if not stanza_source.slides:
            raise SlideGenerationError(f"Stanza template has no slides: {self._stanza_template}")

        stanza_slide_source = stanza_source.slides[0]
        stanza_layout = self._any_layout(presentation)

        for stanza in lyrics.stanzas:
            new_slide = self._duplicate_slide(presentation, stanza_slide_source, stanza_layout)
            self._fill_placeholder_text(new_slide, stanza.text)

        presentation.save(str(combined_path))
        return combined_path

    def _any_layout(self, presentation: Presentation):
        """Returns any existing slide layout from the presentation. The layout choice
        doesn't affect the final slide: _duplicate_slide immediately strips every
        placeholder shape the layout would contribute and replaces the slide's
        content with a deep copy of the source slide's own shapes."""
        return presentation.slide_masters[0].slide_layouts[0]

    def _duplicate_slide(self, presentation: Presentation, source_slide: Slide, layout) -> Slide:
        new_slide = presentation.slides.add_slide(layout)

        for shape in list(new_slide.shapes):
            shape._element.getparent().remove(shape._element)

        for shape in source_slide.shapes:
            new_element = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.append(new_element)

        self._remap_images(source_slide, new_slide)
        return new_slide

    def _remap_images(self, source_slide: Slide, new_slide: Slide) -> None:
        """Re-embeds every image referenced anywhere in new_slide's shape tree,
        including pictures nested inside (possibly nested) groups, resolving the
        original bytes from source_slide's relationships."""
        blip_tag = "{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
        r_embed_attr = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"

        for blip in new_slide.shapes._spTree.iter(blip_tag):
            old_r_id = blip.get(r_embed_attr)
            if not old_r_id:
                continue

            try:
                image_part = source_slide.part.related_part(old_r_id)
            except KeyError:
                continue

            _, new_r_id = new_slide.part.get_or_add_image_part(io.BytesIO(image_part.blob))
            blip.set(r_embed_attr, new_r_id)

    def _fill_placeholder_text(self, slide: Slide, text: str) -> None:
        for shape in self._iter_text_shapes(slide.shapes):
            self._set_text_preserving_format(shape.text_frame, text)
            return

        raise SlideGenerationError("Template slide has no text placeholder to fill.")

    def _set_text_preserving_format(self, text_frame, text: str) -> None:
        """Replaces a text frame's content line-by-line, re-applying the font (name,
        size, bold, italic, color) of the template's original first run to every new
        run. TextFrame.text = ... would otherwise silently discard that formatting."""
        template_font = self._first_run_font(text_frame)
        alignment = text_frame.paragraphs[0].alignment if text_frame.paragraphs else None

        lines = text.split("\n") or [""]

        text_frame.text = lines[0]
        first_paragraph = text_frame.paragraphs[0]
        if alignment is not None:
            first_paragraph.alignment = alignment
        self._apply_font(first_paragraph.runs[0], template_font)

        for line in lines[1:]:
            paragraph = text_frame.add_paragraph()
            if alignment is not None:
                paragraph.alignment = alignment
            run = paragraph.add_run()
            run.text = line
            self._apply_font(run, template_font)

    def _first_run_font(self, text_frame):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                return run.font
        return None

    def _apply_font(self, run, template_font) -> None:
        if template_font is None:
            return

        font = run.font
        font.name = template_font.name
        font.size = template_font.size
        font.bold = template_font.bold
        font.italic = template_font.italic
        if template_font.color and template_font.color.type == MSO_COLOR_TYPE.RGB:
            font.color.rgb = template_font.color.rgb

    def _iter_text_shapes(self, shapes):
        """Depth-first walk of a shape tree, yielding text-frame shapes. Recurses into
        GROUP shapes since real-world templates often nest title/subtitle text boxes
        inside one or more grouped shapes."""
        for shape in shapes:
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                yield from self._iter_text_shapes(shape.shapes)
            elif shape.has_text_frame:
                yield shape


class PowerPointPngExporter:
    """Exports every slide of a .pptx as a numbered PNG using PowerPoint COM automation.

    Windows-only: requires a local PowerPoint installation. This step cannot be
    exercised in a headless/non-Windows environment; verify manually with real
    templates and PowerPoint installed.
    """

    def export(self, pptx_path: Path, output_dir: Path) -> list[Path]:
        import win32com.client  # noqa: PLC0415 (Windows-only, optional dependency)

        output_dir.mkdir(parents=True, exist_ok=True)

        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(
            str(pptx_path.resolve()), WithWindow=False
        )
        try:
            with tempfile.TemporaryDirectory(prefix="hymn-studio-pptx-") as temp_dir:
                presentation.SaveAs(temp_dir, 18)  # ppSaveAsPNG = 18, exports one PNG per slide

                exported = sorted(Path(temp_dir).glob("*.PNG")) or sorted(
                    Path(temp_dir).glob("*.png")
                )
                paths: list[Path] = []
                for index, source in enumerate(exported, start=1):
                    destination = output_dir / f"{index:04d}_slide.png"
                    shutil.copyfile(source, destination)
                    paths.append(destination)
                return paths
        finally:
            presentation.Close()
            powerpoint.Quit()
